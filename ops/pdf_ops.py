"""Core PDF operations using pypdf and reportlab."""
import os
import io
import shutil
import tempfile
from pathlib import Path


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _load_pypdf():
    try:
        from pypdf import PdfReader, PdfWriter
        return PdfReader, PdfWriter
    except ImportError:
        raise ImportError("pypdf non installato. Esegui: pip install pypdf")


def _load_reportlab():
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        return canvas, A4
    except ImportError:
        raise ImportError("reportlab non installato. Esegui: pip install reportlab")


# ---------------------------------------------------------------------------
# ORGANIZE
# ---------------------------------------------------------------------------

def merge_pdfs(input_paths: list, output_path: str) -> str:
    PdfReader, PdfWriter = _load_pypdf()
    writer = PdfWriter()
    for path in input_paths:
        reader = PdfReader(path)
        for page in reader.pages:
            writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def split_pdf(input_path: str, output_dir: str, pages_per_file: int = 1) -> list:
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    total = len(reader.pages)
    stem = Path(input_path).stem
    outputs = []
    for start in range(0, total, pages_per_file):
        writer = PdfWriter()
        chunk = reader.pages[start:start + pages_per_file]
        for page in chunk:
            writer.add_page(page)
        out = os.path.join(output_dir, f"{stem}_part{start // pages_per_file + 1}.pdf")
        with open(out, "wb") as f:
            writer.write(f)
        outputs.append(out)
    return outputs


def remove_pages(input_path: str, output_path: str, pages_to_remove: list) -> str:
    """pages_to_remove: 1-based page numbers."""
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    remove_set = set(p - 1 for p in pages_to_remove)
    for i, page in enumerate(reader.pages):
        if i not in remove_set:
            writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def extract_pages(input_path: str, output_path: str, pages_to_extract: list) -> str:
    """pages_to_extract: 1-based page numbers."""
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for p in pages_to_extract:
        writer.add_page(reader.pages[p - 1])
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def reorder_pages(input_path: str, output_path: str, new_order: list) -> str:
    """new_order: list of 1-based page numbers in the desired order."""
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for p in new_order:
        writer.add_page(reader.pages[p - 1])
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def images_to_pdf(image_paths: list, output_path: str) -> str:
    """Convert one or more image files to a single PDF (scan to PDF / jpg to PDF)."""
    try:
        import img2pdf
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(image_paths))
        return output_path
    except ImportError:
        pass
    # fallback via Pillow
    from PIL import Image
    images = []
    for p in image_paths:
        img = Image.open(p).convert("RGB")
        images.append(img)
    if not images:
        raise ValueError("Nessuna immagine fornita.")
    images[0].save(output_path, save_all=True, append_images=images[1:])
    return output_path


# ---------------------------------------------------------------------------
# OPTIMIZE
# ---------------------------------------------------------------------------

def compress_pdf(input_path: str, output_path: str) -> str:
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def repair_pdf(input_path: str, output_path: str) -> str:
    """Basic repair: read and rewrite via pypdf (fixes many structural issues)."""
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path, strict=False)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def ocr_pdf(input_path: str, output_path: str, lang: str = "ita+eng") -> str:
    try:
        from pdf2image import convert_from_path
        import pytesseract
        from PIL import Image
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.pagesizes import letter
    except ImportError as e:
        raise ImportError(f"Dipendenza mancante per OCR: {e}")

    pages = convert_from_path(input_path, dpi=200)
    pdf_pages = []
    tmpdir = tempfile.mkdtemp()
    try:
        for i, page_img in enumerate(pages):
            pdf_bytes = pytesseract.image_to_pdf_or_hocr(page_img, extension="pdf", lang=lang)
            p = os.path.join(tmpdir, f"page_{i}.pdf")
            with open(p, "wb") as f:
                f.write(pdf_bytes)
            pdf_pages.append(p)
        merge_pdfs(pdf_pages, output_path)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    return output_path


# ---------------------------------------------------------------------------
# CONVERT TO PDF
# ---------------------------------------------------------------------------

def word_to_pdf(input_path: str, output_path: str) -> str:
    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(input_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)
        doc.Close()
        word.Quit()
        return output_path
    except ImportError:
        pass
    # fallback: python-docx → reportlab (plain text)
    from docx import Document as DocxDocument
    canvas_mod, A4 = _load_reportlab()
    doc = DocxDocument(input_path)
    c = canvas_mod.Canvas(output_path, pagesize=A4)
    width, height = A4
    y = height - 50
    for para in doc.paragraphs:
        text = para.text
        if not text.strip():
            y -= 14
        else:
            c.setFont("Helvetica", 11)
            c.drawString(50, y, text[:120])
            y -= 16
        if y < 50:
            c.showPage()
            y = height - 50
    c.save()
    return output_path


def pptx_to_pdf(input_path: str, output_path: str) -> str:
    try:
        import win32com.client
        pp = win32com.client.Dispatch("PowerPoint.Application")
        pp.Visible = True
        prs = pp.Presentations.Open(os.path.abspath(input_path))
        prs.SaveAs(os.path.abspath(output_path), 32)  # 32 = ppSaveAsPDF
        prs.Close()
        pp.Quit()
        return output_path
    except ImportError:
        raise ImportError(
            "Conversione PPTX→PDF richiede Microsoft PowerPoint installato (win32com).\n"
            "Installa pywin32: pip install pywin32"
        )


def excel_to_pdf(input_path: str, output_path: str) -> str:
    try:
        import win32com.client
        xl = win32com.client.Dispatch("Excel.Application")
        xl.Visible = False
        wb = xl.Workbooks.Open(os.path.abspath(input_path))
        wb.ExportAsFixedFormat(0, os.path.abspath(output_path))
        wb.Close(False)
        xl.Quit()
        return output_path
    except ImportError:
        raise ImportError(
            "Conversione EXCEL→PDF richiede Microsoft Excel installato (win32com).\n"
            "Installa pywin32: pip install pywin32"
        )


def html_to_pdf(input_path: str, output_path: str) -> str:
    try:
        import weasyprint
        weasyprint.HTML(filename=input_path).write_pdf(output_path)
        return output_path
    except ImportError:
        pass
    try:
        import pdfkit
        pdfkit.from_file(input_path, output_path)
        return output_path
    except ImportError:
        raise ImportError(
            "Conversione HTML→PDF richiede weasyprint o pdfkit.\n"
            "pip install weasyprint  oppure  pip install pdfkit"
        )


# ---------------------------------------------------------------------------
# CONVERT FROM PDF
# ---------------------------------------------------------------------------

def pdf_to_images(input_path: str, output_dir: str, fmt: str = "JPEG", dpi: int = 150) -> list:
    try:
        from pdf2image import convert_from_path
    except ImportError:
        raise ImportError("pdf2image non installato. Esegui: pip install pdf2image\nRichiede anche Poppler nel PATH.")
    pages = convert_from_path(input_path, dpi=dpi)
    stem = Path(input_path).stem
    ext = "jpg" if fmt.upper() == "JPEG" else fmt.lower()
    outputs = []
    for i, img in enumerate(pages):
        out = os.path.join(output_dir, f"{stem}_page{i + 1}.{ext}")
        img.save(out, fmt)
        outputs.append(out)
    return outputs


def pdf_to_word(input_path: str, output_path: str) -> str:
    try:
        from pdf2docx import Converter
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        return output_path
    except ImportError:
        raise ImportError("pdf2docx non installato. Esegui: pip install pdf2docx")


def pdf_to_excel(input_path: str, output_path: str) -> str:
    try:
        import pdfplumber
        import openpyxl
    except ImportError:
        raise ImportError("pdfplumber/openpyxl non installati. Esegui: pip install pdfplumber openpyxl")
    wb = openpyxl.Workbook()
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages):
            ws = wb.create_sheet(title=f"Page{page_num + 1}")
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    for row in table:
                        ws.append([cell or "" for cell in row])
            else:
                text = page.extract_text() or ""
                for line in text.split("\n"):
                    ws.append([line])
    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]
    wb.save(output_path)
    return output_path


def pdf_to_pdfa(input_path: str, output_path: str) -> str:
    """Basic PDF/A-1b conversion (metadata tagging via pypdf + reportlab)."""
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.add_metadata({
        "/Title": Path(input_path).stem,
        "/Creator": "Universal PDF Converter",
    })
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def pdf_to_pptx(input_path: str, output_path: str) -> str:
    """Convert each PDF page to an image slide in PPTX."""
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("Richiede: pip install pdf2image python-pptx\nE Poppler nel PATH.")
    pages = convert_from_path(input_path, dpi=150)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    tmpdir = tempfile.mkdtemp()
    try:
        for i, img in enumerate(pages):
            img_path = os.path.join(tmpdir, f"slide_{i}.jpg")
            img.save(img_path, "JPEG")
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
        prs.save(output_path)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    return output_path


# ---------------------------------------------------------------------------
# EDIT
# ---------------------------------------------------------------------------

def rotate_pdf(input_path: str, output_path: str, angle: int, pages: list = None) -> str:
    """angle: 90, 180, 270. pages: 1-based list or None for all."""
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    target = set(p - 1 for p in pages) if pages else None
    for i, page in enumerate(reader.pages):
        if target is None or i in target:
            page.rotate(angle)
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def add_page_numbers(input_path: str, output_path: str,
                     position: str = "bottom-center",
                     font_size: int = 10,
                     start_num: int = 1) -> str:
    PdfReader, PdfWriter = _load_pypdf()
    canvas_mod, _ = _load_reportlab()
    reader = PdfReader(input_path)
    writer = PdfWriter()

    for i, page in enumerate(reader.pages):
        box = page.mediabox
        w = float(box.width)
        h = float(box.height)
        packet = io.BytesIO()
        c = canvas_mod.Canvas(packet, pagesize=(w, h))
        c.setFont("Helvetica", font_size)
        num_text = str(i + start_num)
        if position == "bottom-center":
            c.drawCentredString(w / 2, 20, num_text)
        elif position == "bottom-right":
            c.drawRightString(w - 20, 20, num_text)
        elif position == "bottom-left":
            c.drawString(20, 20, num_text)
        elif position == "top-center":
            c.drawCentredString(w / 2, h - 20, num_text)
        c.save()
        packet.seek(0)
        overlay_reader = PdfReader(packet)
        page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def add_watermark(input_path: str, output_path: str,
                  text: str = "CONFIDENTIAL",
                  opacity: float = 0.3,
                  angle: float = 45,
                  font_size: int = 48) -> str:
    PdfReader, PdfWriter = _load_pypdf()
    canvas_mod, _ = _load_reportlab()
    from reportlab.lib.colors import Color

    reader = PdfReader(input_path)
    writer = PdfWriter()

    # Build watermark page using first page dimensions
    first = reader.pages[0]
    box = first.mediabox
    w, h = float(box.width), float(box.height)
    packet = io.BytesIO()
    c = canvas_mod.Canvas(packet, pagesize=(w, h))
    c.saveState()
    c.setFillColor(Color(0, 0, 0, alpha=opacity))
    c.setFont("Helvetica-Bold", font_size)
    c.translate(w / 2, h / 2)
    c.rotate(angle)
    c.drawCentredString(0, 0, text)
    c.restoreState()
    c.save()
    packet.seek(0)
    wm_page = PdfReader(packet).pages[0]

    for page in reader.pages:
        page.merge_page(wm_page)
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def crop_pdf(input_path: str, output_path: str,
             left: float, bottom: float, right: float, top: float) -> str:
    """Crop in points (1 inch = 72 pts). Values from page edges."""
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        box = page.mediabox
        pw = float(box.width)
        ph = float(box.height)
        page.mediabox.lower_left  = (left, bottom)
        page.mediabox.upper_right = (pw - right, ph - top)
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def extract_form_data(input_path: str) -> dict:
    PdfReader, _ = _load_pypdf()
    reader = PdfReader(input_path)
    fields = reader.get_fields()
    return {k: v.get("/V", "") for k, v in (fields or {}).items()}


# ---------------------------------------------------------------------------
# SECURITY
# ---------------------------------------------------------------------------

def protect_pdf(input_path: str, output_path: str,
                user_password: str, owner_password: str = None) -> str:
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_password=user_password,
                   owner_password=owner_password or user_password,
                   use_128bit=True)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def unlock_pdf(input_path: str, output_path: str, password: str) -> str:
    PdfReader, PdfWriter = _load_pypdf()
    reader = PdfReader(input_path)
    if reader.is_encrypted:
        if not reader.decrypt(password):
            raise ValueError("Password errata.")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def redact_pdf(input_path: str, output_path: str, search_text: str) -> str:
    """Black-box redaction: draw black rectangles over matching text."""
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("pdfplumber non installato. Esegui: pip install pdfplumber")
    PdfReader, PdfWriter = _load_pypdf()
    canvas_mod, _ = _load_reportlab()
    from reportlab.lib.colors import black

    reader = PdfReader(input_path)
    writer = PdfWriter()

    with pdfplumber.open(input_path) as pdf:
        for page_idx, (plumber_page, pypdf_page) in enumerate(zip(pdf.pages, reader.pages)):
            box = pypdf_page.mediabox
            pw = float(box.width)
            ph = float(box.height)
            words = plumber_page.extract_words()
            matches = [w for w in words if search_text.lower() in w["text"].lower()]
            if matches:
                packet = io.BytesIO()
                c = canvas_mod.Canvas(packet, pagesize=(pw, ph))
                c.setFillColor(black)
                for m in matches:
                    x0 = m["x0"]
                    y0 = ph - m["bottom"]
                    x1 = m["x1"]
                    y1 = ph - m["top"]
                    c.rect(x0, y0, x1 - x0, y1 - y0, fill=1, stroke=0)
                c.save()
                packet.seek(0)
                overlay = PdfReader(packet).pages[0]
                pypdf_page.merge_page(overlay)
            writer.add_page(pypdf_page)

    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path


def compare_pdfs(path1: str, path2: str) -> dict:
    """Return a dict with page counts and differing page indices (text-based)."""
    PdfReader, _ = _load_pypdf()
    r1 = PdfReader(path1)
    r2 = PdfReader(path2)
    n1, n2 = len(r1.pages), len(r2.pages)
    diffs = []
    for i in range(min(n1, n2)):
        t1 = r1.pages[i].extract_text() or ""
        t2 = r2.pages[i].extract_text() or ""
        if t1.strip() != t2.strip():
            diffs.append(i + 1)
    return {
        "pages_a": n1,
        "pages_b": n2,
        "differing_pages": diffs,
        "identical": n1 == n2 and not diffs,
    }


def sign_pdf(input_path: str, output_path: str, sig_image_path: str,
             page: int = 1, x: float = 50, y: float = 50,
             width: float = 150, height: float = 50) -> str:
    """Stamp a signature image onto the specified page (visual signature)."""
    PdfReader, PdfWriter = _load_pypdf()
    canvas_mod, _ = _load_reportlab()
    reader = PdfReader(input_path)
    writer = PdfWriter()
    target_idx = page - 1
    for i, pdf_page in enumerate(reader.pages):
        if i == target_idx:
            box = pdf_page.mediabox
            pw = float(box.width)
            ph = float(box.height)
            packet = io.BytesIO()
            c = canvas_mod.Canvas(packet, pagesize=(pw, ph))
            c.drawImage(sig_image_path, x, y, width=width, height=height, mask="auto")
            c.save()
            packet.seek(0)
            overlay = PdfReader(packet).pages[0]
            pdf_page.merge_page(overlay)
        writer.add_page(pdf_page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path
